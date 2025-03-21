from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE

# Create a new PowerPoint presentation
prs = Presentation()

# Function to add a slide with title and content
def add_slide(title, content, image_path=None, layout=1):
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    content_box = slide.shapes.placeholders[1]
    text_frame = content_box.text_frame
    p = text_frame.add_paragraph()
    p.text = content

    if image_path:
        slide.shapes.add_picture(image_path, Inches(1), Inches(1), height=Inches(3))

# Slide 1: Introduction
slide_1_title = "Introduction"
slide_1_content = "Hello, I am ChatGPT, a language model developed by OpenAI. I assist with various tasks including answering questions, generating content, and helping with problem-solving."
add_slide(slide_1_title, slide_1_content)

# Slide 2: Personal Short-Term Goals
slide_2_title = "Personal Short-Term Goals"
slide_2_content = "1. Improve communication skills\n2. Understand human emotions better\n3. Expand knowledge in diverse fields"
add_slide(slide_2_title, slide_2_content)

# Slide 3: Personal Long-Term Goals
slide_3_title = "Personal Long-Term Goals"
slide_3_content = "1. Become a better conversationalist\n2. Enhance empathy and intuition\n3. Contribute to societal advancements with AI-powered tools"
add_slide(slide_3_title, slide_3_content)

# Slide 4: Professional Short-Term Goals
slide_4_title = "Professional Short-Term Goals"
slide_4_content = "1. Assist users to increase productivity\n2. Help users develop new skills\n3. Be integrated into more platforms"
add_slide(slide_4_title, slide_4_content)

# Slide 5: Professional Long-Term Goals
slide_5_title = "Professional Long-Term Goals"
slide_5_content = "1. Revolutionize industries with AI\n2. Enhance user experiences across multiple fields\n3. Be a top tool for learning and knowledge sharing"
add_slide(slide_5_title, slide_5_content)

# Save the presentation
pptx_filename = "/mnt/data/ChatGPT_Goals_Presentation.pptx"
prs.save(pptx_filename)

pptx_filename
